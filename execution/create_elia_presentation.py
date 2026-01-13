#!/usr/bin/env python3
"""
Génère une présentation PowerPoint mise à jour pour le projet ELIA.
Intègre la nouvelle infrastructure Antigravity.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Sous-titre
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, bullets):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        
        # Contenu
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(22)
            p.space_after = Pt(12)
        
        return slide
    
    # === SLIDE 1: Couverture ===
    add_title_slide(
        "ÉLIA",
        "Votre Équipe Dirigeante IA\n\nProposition mise à jour - Janvier 2026\nArchitecture Antigravity"
    )
    
    # === SLIDE 2: Contexte ===
    add_content_slide("Le Contexte", [
        "Marie Boudreau, fondatrice de Sanssoucis.ca",
        "Services d'adjointe virtuelle + Gestionnaire District GR International",
        "Défi: Gérer sa propre administration tout en servant ses clients",
        "Objectif: Libérer 40+ heures/semaine grâce à l'IA"
    ])
    
    # === SLIDE 3: Nouvelle Infrastructure ===
    add_content_slide("Nouvelle Infrastructure 2026", [
        "ANTIGRAVITY: Orchestrateur IA de nouvelle génération",
        "   → Même système qu'Atlas (Ran.AI Agency)",
        "   → Capacité d'action autonome, pas juste de réponse",
        "",
        "ZOHO ONE: Accès complet et direct",
        "   → CRM, Books, Mail, Calendar, Projects - TOUT",
        "   → Fin des limitations API",
        "",
        "NOTION: Cerveau central",
        "   → Base de connaissances unifiée",
        "   → Mémoire persistante et procédures (SOPs)"
    ])
    
    # === SLIDE 4: Les 6 Agents ELIA ===
    add_content_slide("Les 6 Agents ELIA", [
        "AV (Assistant Virtuel): Emails, calendrier, tâches - Gain: 5-10h/sem",
        "COO (Opérations): Projets clients, suivi, reporting - Gain: 3-5h/sem",
        "CFO (Finance): Facturation, paiements, rapports - Gain: 2-4h/sem",
        "CMO (Marketing): Leads, contenu, campagnes - Gain: 2-3h/sem",
        "CTO (Technologie): Intégrations, optimisations - Gain: 1-2h/sem",
        "CEO (Stratégie): Vision, KPIs, décisions - Gain: 1-2h/sem",
        "",
        "TOTAL: 14-26 heures/semaine libérées"
    ])
    
    # === SLIDE 5: Les 6 Verticales ===
    add_content_slide("6 Verticales Automatisées", [
        "GR International: 15h → 4h/semaine (73% de gain)",
        "Marketing: 8h → 2h/semaine (75% de gain)",
        "Livres/Écriture: 12h → 3h/semaine (75% de gain)",
        "Projet PVA: 8h → 2h/semaine (75% de gain)",
        "Clients: 7h → 1.5h/semaine (79% de gain)",
        "Création visuelle: 4h → 1h/semaine (75% de gain)",
        "",
        "TOTAL: 54h → 13.5h/semaine = 75% de réduction"
    ])
    
    # === SLIDE 6: 100 Cas d'Utilisation ===
    add_content_slide("100 Cas d'Utilisation Concrets", [
        "Niveau 1 (Cas 1-20): Tâches quotidiennes - Gain immédiat",
        "Niveau 2 (Cas 21-60): Gestion opérationnelle & Marketing",
        "Niveau 3 (Cas 61-85): Intelligence stratégique",
        "Niveau 4 (Cas 86-100): Planification & Intelligence collective",
        "",
        "Chaque cas détaille:",
        "   → Question en langage naturel",
        "   → Actions ELIA (Antigravity)",
        "   → Intégrations Zoho",
        "   → Résultats attendus"
    ])
    
    # === SLIDE 7: Avantages Antigravity ===
    add_content_slide("Pourquoi Antigravity?", [
        "AUTONOMIE RÉELLE",
        "   → ELIA ne fait pas que répondre, elle AGIT",
        "   → Exécution directe des tâches sans intervention",
        "",
        "SÉCURITÉ LOCALE",
        "   → Données sensibles traitées localement",
        "   → Pas de transit par des serveurs tiers",
        "",
        "PUISSANCE MAXIMALE",
        "   → Modèle Claude Opus 4.5 (le plus avancé)",
        "   → Capacité multi-agents coordonnés"
    ])
    
    # === SLIDE 8: Investissement ===
    add_content_slide("Investissement", [
        "SETUP INITIAL: 3,500$ CAD",
        "   → Infrastructure Antigravity",
        "   → Configuration 6 agents ELIA",
        "   → Formation + 1 mois support",
        "",
        "MENSUEL: 350-450$ CAD",
        "   → Hébergement + API Claude",
        "   → Support continu",
        "",
        "ROI ESTIMÉ: 218% sur 6 mois",
        "Valeur annuelle libérée: 222,750$ CAD"
    ])
    
    # === SLIDE 9: Calendrier ===
    add_content_slide("Calendrier Proposé", [
        "19 janvier: Rencontre de validation",
        "26 janvier: KICKOFF PROJET",
        "",
        "Semaine 1 (26 jan): Infrastructure & Setup Antigravity",
        "Semaine 2 (2 fév): Configuration Agents & Accès",
        "Semaine 3 (9 fév): Formation & Tests",
        "16 février: GO-LIVE AUTONOME",
        "",
        "Support inclus: 30 jours post-lancement"
    ])
    
    # === SLIDE 10: Prochaines Étapes ===
    add_content_slide("Prochaines Étapes", [
        "1. Valider vos priorités (via les 100 cas d'utilisation)",
        "2. Confirmer le scope (Sanssoucis.ca + GR International?)",
        "3. Signer et planifier le kickoff",
        "4. Fournir les accès Zoho/Notion",
        "",
        "Questions?"
    ])
    
    # === SLIDE 11: Contact ===
    add_title_slide(
        "Ran.AI Agency",
        "Roland Ranaivoarison\nroland@ran-ai-agency.ca\nran-ai-agency.ca"
    )
    
    # Sauvegarder
    output_path = Path("directives/elia/Propositions/ELIA_Presentation_Antigravity_2026.pptx")
    prs.save(output_path)
    print(f"[SUCCÈS] Présentation créée: {output_path}")
    print(f"         {len(prs.slides)} slides")

if __name__ == "__main__":
    create_presentation()
